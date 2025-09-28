from typing import Tuple, List
from pptx import Presentation
from backend.app.models.schemas import Paragraph
from backend.app.utils.text_clean import normalize_text

def parse_pptx(path: str) -> Tuple[str, List[Paragraph]]:
    prs = Presentation(path)
    items = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    texts.append(t)
        if texts:
            items.append("\n".join(texts))
    full_text = normalize_text("\n\n".join(items))

    paragraphs: List[Paragraph] = []
    cursor = 0
    for i, t in enumerate(items):
        t_norm = normalize_text(t)
        start = full_text.find(t_norm, cursor)
        if start < 0:
            start = full_text.find(t_norm)
        end = start + len(t_norm) if start >= 0 else cursor
        paragraphs.append(Paragraph(index=i, text=t_norm, start_char=max(0,start), end_char=max(end,start)))
        cursor = end
    return full_text, paragraphs
