"""Utilities for splitting different document types into page-wise text blocks."""

from __future__ import annotations

from dataclasses import dataclass
from typing import List

from backend.app.utils.text_clean import normalize_text


@dataclass
class PageContent:
    page_number: int
    text: str


def _parse_pdf(path: str) -> List[PageContent]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    pages: List[PageContent] = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        pages.append(PageContent(page_number=idx, text=normalize_text(text)))
    return pages


def _parse_pptx(path: str) -> List[PageContent]:
    from pptx import Presentation

    prs = Presentation(path)
    pages: List[PageContent] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    texts.append(text)
        pages.append(PageContent(page_number=idx, text=normalize_text("\n".join(texts))))
    return pages


def _parse_docx(path: str) -> List[PageContent]:
    from docx import Document

    doc = Document(path)
    buffer: List[str] = []
    pages: List[PageContent] = []
    char_budget = 0
    page_number = 1
    for para in doc.paragraphs:
        text = normalize_text(para.text.strip())
        if not text:
            continue
        buffer.append(text)
        char_budget += len(text)
        if char_budget >= 1200:
            pages.append(PageContent(page_number=page_number, text="\n".join(buffer)))
            page_number += 1
            buffer = []
            char_budget = 0
    if buffer:
        pages.append(PageContent(page_number=page_number, text="\n".join(buffer)))
    return pages or [PageContent(page_number=1, text="")]


def _parse_plain_text(path: str) -> List[PageContent]:
    with open(path, "r", encoding="utf-8", errors="ignore") as handle:
        content = normalize_text(handle.read())
    chunks: List[PageContent] = []
    words = content.splitlines()
    buffer: List[str] = []
    page_number = 1
    for line in words:
        if len("\n".join(buffer)) > 1500:
            chunks.append(PageContent(page_number=page_number, text="\n".join(buffer)))
            page_number += 1
            buffer = []
        buffer.append(line)
    if buffer:
        chunks.append(PageContent(page_number=page_number, text="\n".join(buffer)))
    return chunks or [PageContent(page_number=1, text="")]


def parse_pages(path: str, extension: str) -> List[PageContent]:
    ext = extension.lower()
    if ext == ".pdf":
        pages = _parse_pdf(path)
    elif ext in {".ppt", ".pptx"}:
        pages = _parse_pptx(path)
    elif ext in {".doc", ".docx"}:
        pages = _parse_docx(path)
    elif ext in {".md", ".txt"}:
        pages = _parse_plain_text(path)
    else:
        raise ValueError(f"不支援的副檔名: {ext}")

    return pages
